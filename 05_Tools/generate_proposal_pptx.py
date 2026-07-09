#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
bravesoft 提案書 PPTX 生成ツール（Step 08）

見積もり JSON（Excel 生成と同じ入力 = 1 つの正本）から、bravesoft ブランドの
「明るく・図解の多い・一目で分かる」提案書デッキ(.pptx)を生成する。

スライド構成（データがあるものだけ出す）:
  01 表紙                         宛先 / 案件名 / 御提案書
  02 本ご提案のポイント            総工数・期間・提示額・主要機能数の 4 スタット
  03 ── セクション 01: 安心の根拠
  04 このお見積もりが「安心」な3つの理由（積み上げ/抜け漏れゼロ/前提明確）
  05 費用構成（フェーズ別ドーナツ）
  06 ── セクション 02: ご提案スコープ
  07 スコープ・主要機能（Must / Want の2カラム）
  08 ── セクション 03: お見積もり
  09 進め方（単一案=フェーズ流れ / 複数案=案比較カード）
  10 お見積もり内訳（フェーズ表＋合計カード）
  11 工数の内訳（フェーズ別 横棒）
  12 ── セクション 04: 概算スケジュール
  13 概算スケジュール（フェーズ別ガント）
  14 まとめ（6つの安心ポイント＋次のステップ）
  15 クロージング（ありがとうございました）

使い方:
  python3 generate_proposal_pptx.py --input sample_input.json --out 提案書.pptx

※ 顧客名・案件名は入力 JSON の meta に従う（テンプレは ○○ プレースホルダ）。
※ この提案書は「提出用」= 社内原価・利益率は載せない（Rule 10）。
"""

import argparse
import io
import json
import math
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- ブランド
BLUE = RGBColor(0x00, 0x66, 0xFF)
SKY = RGBColor(0x3F, 0x9C, 0xFF)
NAVY = RGBColor(0x25, 0x3A, 0x58)
INK = RGBColor(0x09, 0x0A, 0x0A)
TINT = RGBColor(0xE5, 0xF0, 0xFF)
SURFACE = RGBColor(0xF7, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6C, 0x70, 0x72)
HAIR = RGBColor(0xE3, 0xE5, 0xE5)
GREEN = RGBColor(0x00, 0xA8, 0x6B)
FONT = "Noto Sans JP"
MONO = "Roboto Mono"
# ドーナツ/棒の配色（青系グラデ）
PALETTE = [BLUE, SKY, NAVY, RGBColor(0x6F, 0xB3, 0xFF), RGBColor(0x1E, 0x5B, 0xC6),
           RGBColor(0x9C, 0xC9, 0xFF), RGBColor(0x2C, 0x4A, 0x7A), RGBColor(0xC4, 0xE0, 0xFF)]

EMU_IN = 914400
SLIDE_W = 13.333
SLIDE_H = 7.5
LOGO = os.path.join(os.path.dirname(__file__), "..", "assets", "bravesoft-logo.png")
LOGO_AR = 1858 / 801  # 幅/高さ


# ---------------------------------------------------------------- 低レベル helper
def _set_radius(shape, inches):
    """角丸四角形の丸み半径を絶対量で固定（サイズに対する既定比だと大きくなりすぎる）。"""
    try:
        target = Emu(int(Inches(inches)))
        ss = min(shape.width, shape.height)
        adj = max(0.0, min(0.5, float(target) / float(ss))) if ss else 0.1
        shape.adjustments[0] = adj
    except Exception:
        pass


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, radius=None, shadow=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        _set_radius(s, radius)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if shadow:
        _soft_shadow(s)
    return s


def _soft_shadow(shape):
    """やわらかいドロップシャドウ（カードの浮き）。"""
    sp = shape._element.spPr
    effs = sp.find(qn("a:effectLst"))
    if effs is None:
        effs = sp.makeelement(qn("a:effectLst"), {})
        sp.append(effs)
    sh = effs.makeelement(qn("a:outerShdw"),
                          {"blurRad": "90000", "dist": "38100", "dir": "5400000", "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "253A58"})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": "18000"}))
    sh.append(clr)
    effs.append(sh)


def oval(slide, x, y, d, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.12, wrap=True, shrink=False):
    """runs: str か [(text, {overrides})] のリスト。段落は '\n' で分割。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.vertical_anchor = anchor
    if shrink:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.NONE
    if isinstance(runs, str):
        runs = [(runs, {})]
    # 段落ごとに: 各 run 内の \n を段落分割に展開
    paras = [[]]
    for t, ov in runs:
        parts = t.split("\n")
        for i, part in enumerate(parts):
            if i > 0:
                paras.append([])
            paras[-1].append((part, ov))
    for pi, para in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for t, ov in para:
            r = p.add_run()
            r.text = t
            r.font.name = ov.get("font", font)
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.color.rgb = ov.get("color", color)
            _set_ea_font(r, ov.get("font", font))
    return tb


def _set_ea_font(run, name):
    """日本語（East Asian）フォントも同じ名前に設定。"""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def white_logo_stream():
    """透過ロゴの色を白にした版を作る（濃色背景用）。PIL が無ければ None。"""
    try:
        from PIL import Image
        im = Image.open(LOGO).convert("RGBA")
        px = im.load()
        for yy in range(im.height):
            for xx in range(im.width):
                r, g, b, a = px[xx, yy]
                if a > 0:
                    px[xx, yy] = (255, 255, 255, a)
        buf = io.BytesIO()
        im.save(buf, "PNG")
        buf.seek(0)
        return buf
    except Exception:
        return None


def add_logo(slide, x, y, w, white=False):
    src = white_logo_stream() if white else LOGO
    if src is None:
        return
    try:
        slide.shapes.add_picture(src, Inches(x), Inches(y), width=Inches(w))
    except Exception:
        pass


def yen(n):
    return "¥{:,.0f}".format(round(n))


# ---------------------------------------------------------------- 集計（Excel と同ロジック）
def compute(data):
    p = data.get("params", {})
    rates = {r["key"]: (r.get("internal", 0), r.get("external", 0)) for r in data.get("rates", [])}
    pats = []
    for pat in data.get("patterns", []):
        flat = []
        for it in pat["items"]:
            k = it.get("type", "normal")
            if k == "parent":
                for ch in it["children"]:
                    c = dict(ch); c["kind"] = "leaf"; c.setdefault("phase", it["phase"]); flat.append(c)
            elif k == "pm":
                c = dict(it); c["kind"] = "pm"; flat.append(c)
            else:
                c = dict(it); c["kind"] = "leaf"; flat.append(c)
        phases = []
        for it in flat:
            if it["phase"] not in phases:
                phases.append(it["phase"])
        for it in flat:
            if it["kind"] == "leaf":
                ok = it.get("checked", True) and it.get("role") in rates
                d = it.get("days", 0) or 0
                it["ext"] = round(d * rates[it["role"]][1]) if ok else 0
                it["days_e"] = d if it.get("checked", True) else 0
            else:
                it["ext"] = 0; it["days_e"] = 0
        for ph in phases:
            be = sum(x["ext"] for x in flat if x["phase"] == ph and x["kind"] == "leaf")
            for x in flat:
                if x["phase"] == ph and x["kind"] == "pm":
                    rate = x.get("rate") or p.get("pm_rate", 0.2)
                    x["ext"] = round(be * rate)
        phase_rows = []
        for ph in phases:
            rows = [x for x in flat if x["phase"] == ph]
            phase_rows.append({
                "phase": ph,
                "days": sum(x["days_e"] for x in rows if x["kind"] == "leaf"),
                "ext": sum(x["ext"] for x in rows),
            })
        T_ext = sum(r["ext"] for r in phase_rows)
        T_days = sum(r["days"] for r in phase_rows)
        buf = round(T_ext * (p.get("buffer", 1.2) - 1)) if p.get("apply_buffer") else 0
        ret = round(T_ext * p.get("retainer_rate", 0))
        quote = T_ext + buf + ret - p.get("discount", 0)
        tax = round(quote * p.get("tax", 0.1))
        # Must / Want 機能（leaf・計上のみ）
        musts = [x for x in flat if x["kind"] == "leaf" and x.get("checked", True)
                 and x.get("mwo", "M") == "M" and x.get("category") != "PM・全体管理"]
        wants = [x for x in flat if x["kind"] == "leaf" and x.get("checked", True)
                 and x.get("mwo") == "W"]
        pats.append({
            "name": pat.get("name", ""), "phase_rows": phase_rows,
            "T_ext": T_ext, "T_days": T_days, "buf": buf, "ret": ret,
            "quote": quote, "tax": tax, "incl": quote + tax,
            "musts": musts, "wants": wants,
            "days_per_month": p.get("days_per_month", 20) or 20,
        })
    return pats


# ---------------------------------------------------------------- スライド部品
def new(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def bg(slide, color):
    rect(slide, -0.06, -0.06, SLIDE_W + 0.12, SLIDE_H + 0.12, fill=color)


def kicker_num(slide, x, y, n, on_dark=False):
    text(slide, x, y, 1.2, 0.5, str(n).zfill(2), size=15, bold=True,
         color=(SKY if on_dark else BLUE), font=MONO)


def footer(slide, on_dark=False):
    c = RGBColor(0xC7, 0xD6, 0xEC) if on_dark else MUTED
    text(slide, 0.6, SLIDE_H - 0.42, 9, 0.3,
         "© bravesoft Inc.", size=9, color=c)
    add_logo(slide, SLIDE_W - 1.35, SLIDE_H - 0.5, 0.75, white=on_dark)


def section_divider(prs, num, title, subtitle):
    s = new(prs)
    bg(s, NAVY)
    # 大きな透かし（右下）
    text(s, SLIDE_W - 7.2, SLIDE_H - 2.7, 7.0, 2.2, "bravesoft", size=66, bold=True,
         color=RGBColor(0x2E, 0x46, 0x69), align=PP_ALIGN.RIGHT)
    text(s, 0.95, 2.0, 2, 1.0, num, size=54, bold=True, color=SKY, font=MONO)
    text(s, 1.0, 3.05, 11, 1.2, title, size=38, bold=True, color=WHITE)
    text(s, 1.02, 4.25, 11, 0.8, subtitle, size=15, color=RGBColor(0xC7, 0xD6, 0xEC))
    add_logo(s, 0.95, 0.6, 1.5, white=True)
    return s


def content_head(slide, num, title, subtitle):
    """コンテンツスライドの見出し（番号バッジ＋タイトル＋サブ）。"""
    oval(slide, 0.6, 0.62, 0.42, BLUE)
    text(slide, 0.6, 0.62, 0.42, 0.42, str(num), size=15, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, 1.18, 0.5, 11.4, 0.6, title, size=25, bold=True, color=INK)
    if subtitle:
        text(slide, 1.18, 1.12, 11.4, 0.5, subtitle, size=12.5, color=MUTED)


# ---------------------------------------------------------------- 各スライド
def slide_cover(prs, meta):
    s = new(prs)
    bg(s, WHITE)
    # 右側の明るいブロック（図解モチーフ）
    rect(s, 9.15, 0, SLIDE_W - 9.15 + 0.1, SLIDE_H, fill=TINT)
    oval(s, 10.15, -1.4, 4.2, BLUE)
    o = oval(s, 11.2, 3.9, 3.0, SKY)
    o.fill.fore_color.rgb = SKY
    add_logo(s, 0.9, 0.75, 2.2)
    text(s, 0.92, 2.5, 7.9, 0.5,
         (meta.get("client", "株式会社○○") + " 御中"), size=15, bold=True, color=NAVY)
    text(s, 0.9, 3.05, 8.0, 1.8, meta.get("project", "○○アプリ / システム開発"),
         size=37, bold=True, color=INK, line_spacing=1.1)
    # 御提案書バッジ
    b = rect(s, 0.92, 5.0, 1.9, 0.5, fill=BLUE, radius=0.25)
    text(s, 0.92, 5.0, 1.9, 0.5, "御 提 案 書", size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.95, 5.75, 7.8, 0.6,
         "機能を分解し、工数×単価を積み上げた根拠のあるお見積もりです。",
         size=12.5, color=MUTED)
    date = meta.get("date", "")
    text(s, 0.95, SLIDE_H - 0.7, 8, 0.4,
         (date + "　／　bravesoft株式会社").strip("　／ "), size=11, color=MUTED)
    return s


def slide_points(prs, pat, meta):
    s = new(prs)
    bg(s, WHITE)
    content_head(s, "", "本ご提案のポイント",
                 "勘や相場ではなく、機能を実装粒度まで分解して工数を積み上げています。")
    months = max(1, math.ceil(pat["T_days"] / pat["days_per_month"]))
    stats = [
        ("総工数", "{:,}".format(pat["T_days"]), "人日", BLUE),
        ("想定期間", "約{}".format(months), "ヶ月", SKY),
        ("主要機能", "{}".format(len(pat["musts"]) + len(pat["wants"])), "項目", NAVY),
        ("ご提示額(税込)", "{:,}".format(round(pat["incl"] / 10000)), "万円", BLUE),
    ]
    cw, gap, x0, y0, ch = 2.86, 0.25, 0.72, 1.95, 1.95
    for i, (label, big, unit, col) in enumerate(stats):
        x = x0 + i * (cw + gap)
        rect(s, x, y0, cw, ch, fill=WHITE, line=HAIR, radius=0.12, shadow=True)
        oval(s, x + 0.28, y0 + 0.32, 0.16, col)
        text(s, x + 0.54, y0 + 0.26, cw - 0.7, 0.35, label, size=12, bold=True, color=MUTED)
        text(s, x + 0.28, y0 + 0.72, cw - 0.56, 0.85,
             [(big, {"size": 40, "bold": True, "color": col, "font": MONO}),
              ("  " + unit, {"size": 13, "bold": True, "color": MUTED})],
             anchor=MSO_ANCHOR.MIDDLE)
    # おすすめの進め方 callout
    cy = 4.35
    rect(s, x0, cy, (cw + gap) * 4 - gap, 1.9, fill=TINT, radius=0.14)
    oval(s, x0 + 0.35, cy + 0.32, 0.5, BLUE)
    text(s, x0 + 0.35, cy + 0.32, 0.5, 0.5, "★", size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x0 + 1.05, cy + 0.3, 10.5, 0.5,
         "おすすめの進め方：まず要件定義フェーズから、リスクを抑えてスタート",
         size=15, bold=True, color=NAVY)
    text(s, x0 + 1.05, cy + 0.95, 10.3, 0.8,
         "上流（要件定義）で仕様・スコープ・お見積もりを精緻化 → 本開発へ。"
         "前提が固まった状態で着手するため、手戻り・想定外を最小化します。",
         size=12.5, color=INK, line_spacing=1.3)
    footer(s)
    return s


def slide_reasons(prs):
    s = new(prs)
    bg(s, WHITE)
    content_head(s, "", "このお見積もりが「安心」な3つの理由",
                 "机上の概算ではなく、再現性のある手順で作成しています。")
    cards = [
        ("1", "積み上げの根拠", "機能・作業を項目まで分解し、工数(人日)×単価で 1 円まで積み上げ。"
         "「なぜこの金額か」に必ず答えられます。"),
        ("2", "抜け漏れゼロ", "PM・テスト・リリース・非機能要件まで、フェーズ×カテゴリで"
         "機械的に網羅。属人的な抜けを防ぎます。"),
        ("3", "前提・スコープが明確", "対応範囲・対象外・契約形態・納品物を明記。"
         "「ここから外れたら別途」を最初に握り、後の揉め事を防ぎます。"),
    ]
    cw, gap, x0, y0, ch = 3.95, 0.32, 0.72, 2.0, 3.7
    for i, (n, title, body) in enumerate(cards):
        x = x0 + i * (cw + gap)
        rect(s, x, y0, cw, ch, fill=WHITE, line=HAIR, radius=0.14, shadow=True)
        oval(s, x + 0.45, y0 + 0.5, 0.85, TINT)
        text(s, x + 0.45, y0 + 0.5, 0.85, 0.85, n, size=30, bold=True, color=BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
        text(s, x + 0.45, y0 + 1.6, cw - 0.9, 0.6, title, size=17, bold=True, color=NAVY)
        text(s, x + 0.45, y0 + 2.25, cw - 0.9, 1.3, body, size=12.5, color=INK, line_spacing=1.4)
    footer(s)
    return s


def slide_donut(prs, pat):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    s = new(prs)
    bg(s, WHITE)
    content_head(s, "", "費用構成｜フェーズ別の内訳",
                 "金額がどのフェーズに配分されているか（税抜・ご提示ベース）。")
    rows = [r for r in pat["phase_rows"] if r["ext"] > 0]
    cd = CategoryChartData()
    cd.categories = [r["phase"] for r in rows]
    cd.add_series("費用", tuple(r["ext"] for r in rows))
    gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.7), Inches(1.85),
                            Inches(6.2), Inches(5.0), cd)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = False  # 右側の凡例カードと重複するため非表示
    ser = ch.series[0]
    for i, pt in enumerate(ser.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = PALETTE[i % len(PALETTE)]
    # 右側: 合計と各フェーズ金額の凡例カード
    x = 7.35
    rect(s, x, 1.95, 5.25, 1.15, fill=NAVY, radius=0.14)
    text(s, x + 0.35, 2.12, 4.6, 0.35, "ご提示額（税抜）", size=12, bold=True,
         color=RGBColor(0xC7, 0xD6, 0xEC))
    text(s, x + 0.35, 2.42, 4.6, 0.55, yen(pat["quote"]), size=26, bold=True,
         color=WHITE, font=MONO)
    yy = 3.35
    for i, r in enumerate(rows):
        oval(s, x + 0.05, yy + 0.06, 0.2, PALETTE[i % len(PALETTE)])
        text(s, x + 0.4, yy, 3.2, 0.35, r["phase"], size=12.5, color=INK)
        pct = r["ext"] / pat["T_ext"] * 100 if pat["T_ext"] else 0
        text(s, x + 3.0, yy, 2.25, 0.35,
             [(yen(r["ext"]), {"size": 12.5, "bold": True, "color": NAVY, "font": MONO}),
              ("  {:.0f}%".format(pct), {"size": 11, "color": MUTED})],
             align=PP_ALIGN.RIGHT)
        yy += 0.52
    footer(s)
    return s


def slide_scope(prs, pat):
    s = new(prs)
    bg(s, WHITE)
    content_head(s, "", "ご提案スコープ｜主要機能",
                 "今回のお見積もりに含む機能。Must＝必須、Want＝ご予算に応じて選択可。")
    cols = [("必須（Must）", pat["musts"], BLUE, "この構成で成立する最小スコープ"),
            ("ご希望に応じて（Want）", pat["wants"], SKY, "予算・優先度に合わせて調整できる範囲")]
    cw, gap, x0, y0 = 6.05, 0.35, 0.72, 1.95
    ch = 4.75
    for ci, (title, items, col, note) in enumerate(cols):
        x = x0 + ci * (cw + gap)
        rect(s, x, y0, cw, ch, fill=(WHITE), line=HAIR, radius=0.14, shadow=True)
        head = rect(s, x, y0, cw, 0.78, fill=col, radius=0.14)
        text(s, x + 0.4, y0, cw - 0.8, 0.78, title, size=16, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + 0.4, y0 + 0.9, cw - 0.8, 0.35, note, size=11, color=MUTED)
        yy = y0 + 1.35
        cap = 8
        shown = items[:cap]
        for it in shown:
            oval(s, x + 0.42, yy + 0.07, 0.14, col)
            label = it.get("item", "").replace("　└ ", "").replace("【内訳合計】", "")
            text(s, x + 0.72, yy, cw - 1.15, 0.4, label, size=12.5, color=INK,
                 anchor=MSO_ANCHOR.MIDDLE)
            yy += 0.38
        if len(items) > cap:
            text(s, x + 0.72, yy + 0.02, cw - 1.1, 0.35,
                 "ほか {} 項目".format(len(items) - cap), size=11, bold=True, color=col)
        if not items:
            text(s, x + 0.72, yy, cw - 1.1, 0.4, "（なし）", size=12, color=MUTED)
    footer(s)
    return s


def slide_approach(prs, pats):
    s = new(prs)
    bg(s, WHITE)
    if len(pats) > 1:
        content_head(s, "", "進め方｜複数プランのご提案",
                     "スコープ・段階に応じた複数案。ご予算と優先度に合わせて選べます。")
        cw, gap, x0, y0, ch = 3.95, 0.32, 0.72, 2.05, 3.6
        for i, pat in enumerate(pats[:3]):
            x = x0 + i * (cw + gap)
            col = PALETTE[i % 3]
            rect(s, x, y0, cw, ch, fill=WHITE, line=HAIR, radius=0.14, shadow=True)
            head = rect(s, x, y0, cw, 0.7, fill=col, radius=0.14)
            text(s, x, y0, cw, 0.7, pat["name"] or "案 {}".format(i + 1), size=15, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            text(s, x + 0.4, y0 + 1.0, cw - 0.8, 0.9,
                 [(yen(pat["quote"]), {"size": 25, "bold": True, "color": NAVY, "font": MONO}),
                  ("\n税抜", {"size": 11, "color": MUTED})])
            months = max(1, math.ceil(pat["T_days"] / pat["days_per_month"]))
            text(s, x + 0.4, y0 + 2.05, cw - 0.8, 1.3,
                 [("工数 ", {"size": 12, "color": MUTED}),
                  ("{:,}人日".format(pat["T_days"]), {"size": 13, "bold": True, "color": INK}),
                  ("\n期間 ", {"size": 12, "color": MUTED}),
                  ("約{}ヶ月".format(months), {"size": 13, "bold": True, "color": INK}),
                  ("\n税込 ", {"size": 12, "color": MUTED}),
                  (yen(pat["incl"]), {"size": 13, "bold": True, "color": INK})],
                 line_spacing=1.5)
    else:
        pat = pats[0]
        content_head(s, "", "進め方｜フェーズの流れ",
                     "要件定義から順に進め、各フェーズで成果物を確定していきます。")
        rows = pat["phase_rows"]
        n = len(rows)
        gap = 0.28
        x0 = 0.72
        cw = (SLIDE_W - x0 * 2 - gap * (n - 1)) / n
        y0 = 2.4
        for i, r in enumerate(rows):
            x = x0 + i * (cw + gap)
            rect(s, x, y0, cw, 2.3, fill=WHITE, line=HAIR, radius=0.12, shadow=True)
            oval(s, x + cw / 2 - 0.28, y0 + 0.3, 0.56, TINT)
            text(s, x + cw / 2 - 0.28, y0 + 0.3, 0.56, 0.56, str(i + 1), size=19, bold=True,
                 color=BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
            text(s, x + 0.15, y0 + 1.05, cw - 0.3, 0.6, r["phase"], size=13.5, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)
            text(s, x + 0.15, y0 + 1.62, cw - 0.3, 0.55,
                 "{:,}人日".format(r["days"]), size=12, color=MUTED, align=PP_ALIGN.CENTER)
            if i < n - 1:
                text(s, x + cw - 0.02, y0 + 0.75, gap + 0.04, 0.5, "→", size=16, bold=True,
                     color=SKY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x0, y0 + 2.7, SLIDE_W - x0 * 2, 0.5,
             "※ まず要件定義フェーズで仕様・スコープを固め、リスクを抑えて本開発に着手することをおすすめします。",
             size=11.5, color=MUTED)
    footer(s)
    return s


def slide_breakdown(prs, pat):
    s = new(prs)
    bg(s, WHITE)
    content_head(s, "", "お見積もり内訳",
                 "フェーズ別の工数と金額（税抜）。PM・管理費は各フェーズ小計に含みます。")
    rows = pat["phase_rows"]
    tx, ty, tw = 0.72, 1.95, 8.4
    rh = 0.46
    # ヘッダ
    rect(s, tx, ty, tw, rh, fill=NAVY)
    for cx, cw, label, al in [(tx, 4.4, "フェーズ", PP_ALIGN.LEFT),
                              (tx + 4.4, 1.8, "工数(人日)", PP_ALIGN.RIGHT),
                              (tx + 6.2, 2.2, "金額(税抜)", PP_ALIGN.RIGHT)]:
        pad = 0.25
        text(s, cx + pad, ty, cw - pad * 2, rh, label, size=12.5, bold=True, color=WHITE,
             align=al, anchor=MSO_ANCHOR.MIDDLE)
    yy = ty + rh
    for i, r in enumerate(rows):
        if i % 2 == 1:
            rect(s, tx, yy, tw, rh, fill=SURFACE)
        text(s, tx + 0.25, yy, 4.0, rh, r["phase"], size=12.5, color=INK,
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, tx + 4.4, yy, 1.55, rh, "{:,}".format(r["days"]), size=12.5, color=INK,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
        text(s, tx + 6.2, yy, 1.95, rh, yen(r["ext"]), size=12.5, color=INK,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
        yy += rh
    # 合計行
    rect(s, tx, yy, tw, rh, fill=TINT)
    text(s, tx + 0.25, yy, 4.0, rh, "合計（税抜）", size=13, bold=True, color=NAVY,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, tx + 4.4, yy, 1.55, rh, "{:,}".format(pat["T_days"]), size=13, bold=True,
         color=NAVY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
    text(s, tx + 6.2, yy, 1.95, rh, yen(pat["T_ext"]), size=13, bold=True, color=NAVY,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
    # 右: 合計カード（行を積み上げてから高さを確定＝バッファ非適用でも隙間なし）
    cx, cyy, cwd = 9.5, 1.95, 3.3
    lines = [("小計（税抜）", pat["T_ext"], False)]
    if pat["buf"]:
        lines.append(("リスクバッファ", pat["buf"], False))
    if pat["ret"]:
        lines.append(("体制維持費", pat["ret"], False))
    lines.append(("消費税", pat["tax"], False))
    lines.append(("ご提示額（税込）", pat["incl"], True))
    card_h = 0.4 + sum(0.98 if big else 0.72 for _, _, big in lines) + 0.15
    rect(s, cx, cyy, cwd, card_h, fill=NAVY, radius=0.16, shadow=True)
    yy = 0.38
    for label, val, big in lines:
        text(s, cx + 0.35, cyy + yy, cwd - 0.7, 0.3, label,
             size=(11.5 if big else 11), bold=big, color=RGBColor(0xC7, 0xD6, 0xEC))
        text(s, cx + 0.35, cyy + yy + 0.3, cwd - 0.6, (0.55 if big else 0.4),
             yen(val), size=(22 if big else 15), bold=True,
             color=(SKY if big else WHITE), font=MONO, wrap=False)
        yy += 0.98 if big else 0.72
    text(s, tx, 6.5, tw, 0.5,
         "※ 金額は税抜表記（合計欄を除く）。本前提から外れる場合は別途お見積もりとなります。",
         size=10.5, color=MUTED)
    footer(s)
    return s


def slide_effort_bar(prs, pat):
    s = new(prs)
    bg(s, WHITE)
    content_head(s, "", "工数の内訳｜どこに費用がかかるか",
                 "フェーズ別の工数(人日)。金額の大半は開発フェーズが占めるのが一般的です。")
    rows = sorted([r for r in pat["phase_rows"] if r["days"] > 0],
                  key=lambda r: r["days"], reverse=True)
    mx = max((r["days"] for r in rows), default=1)
    x0, y0 = 0.72, 2.15
    barmax = 8.4
    rh = (4.4 / max(len(rows), 1))
    rh = min(rh, 0.72)
    for i, r in enumerate(rows):
        y = y0 + i * (rh + 0.12)
        text(s, x0, y, 2.3, rh, r["phase"], size=12.5, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
        rect(s, x0 + 2.4, y + rh * 0.15, barmax, rh * 0.7, fill=SURFACE, radius=0.06)
        w = max(0.12, barmax * r["days"] / mx)
        rect(s, x0 + 2.4, y + rh * 0.15, w, rh * 0.7, fill=PALETTE[i % len(PALETTE)], radius=0.06)
        text(s, x0 + 2.4 + w + 0.12, y, 1.6, rh,
             "{:,}人日".format(r["days"]), size=12, bold=True, color=INK,
             anchor=MSO_ANCHOR.MIDDLE, font=MONO)
    footer(s)
    return s


def slide_gantt(prs, pat):
    s = new(prs)
    bg(s, WHITE)
    content_head(s, "", "概算スケジュール",
                 "フェーズを順に進めた場合の目安（並走・前倒しで短縮可能）。")
    rows = pat["phase_rows"]
    dpm = pat["days_per_month"]
    # フェーズ毎の月数（直列・切り上げ）
    segs = []
    cur = 0.0
    for r in rows:
        dur = max(0.5, r["days"] / dpm)
        segs.append((r["phase"], cur, dur))
        cur += dur
    total_m = max(1, math.ceil(cur))
    total_m = min(total_m, 14)
    # グリッド
    gx, gy = 3.3, 2.05
    gw = SLIDE_W - gx - 0.6
    colw = gw / total_m
    rh = min(0.62, 4.4 / max(len(rows), 1))
    # 月ヘッダ
    for m in range(total_m):
        text(s, gx + m * colw, gy - 0.42, colw, 0.35, "{}".format(m + 1) + "ヶ月",
             size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
        ln = rect(s, gx + m * colw, gy, 0.008, rh * len(rows) + 0.2, fill=HAIR)
    for i, (ph, start, dur) in enumerate(segs):
        y = gy + i * rh
        text(s, 0.72, y, 2.4, rh, ph, size=12, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
        bx = gx + start * colw
        bw = max(0.2, dur * colw - 0.06)
        rect(s, bx + 0.03, y + rh * 0.16, bw, rh * 0.62, fill=PALETTE[i % len(PALETTE)],
             radius=0.08)
    # リリース/公開マイルストーン（フェーズに「リリース」が無い場合のみ 1 行追加）
    note_y = gy + len(segs) * rh + 0.15
    if not any("リリース" in ph or "公開" in ph for ph, _, _ in segs):
        ry = gy + len(segs) * rh
        text(s, 0.72, ry, 2.4, rh, "本番公開", size=12, bold=True, color=GREEN,
             anchor=MSO_ANCHOR.MIDDLE)
        rect(s, gx + min(cur, total_m) * colw - colw * 0.5, ry + rh * 0.16,
             colw * 0.45, rh * 0.62, fill=GREEN, radius=0.08)
        note_y = gy + (len(segs) + 1) * rh + 0.15
    text(s, 0.72, note_y, 12, 0.4,
         "総期間の目安：約{}ヶ月（アプリ複数名体制・直列ケース。並走で短縮可能）".format(total_m),
         size=11.5, color=MUTED)
    footer(s)
    return s


def slide_summary(prs):
    s = new(prs)
    bg(s, WHITE)
    content_head(s, "", "まとめ｜このお見積もりが「安心」な理由",
                 "根拠の透明性・過剰のない適正工数・段階的なリスク低減。")
    cards = [
        ("根拠の透明性", "全項目を工数×単価で 1 円まで積み上げ。"),
        ("適正な工数", "相場や勘ではなく、機能の実装粒度から算出。"),
        ("抜け漏れの回避", "PM・テスト・リリース・非機能まで網羅。"),
        ("前提の明確化", "対応範囲・対象外・契約形態を明記。"),
        ("採算の健全性", "バッファ・管理費を見える化し、無理のない体制。"),
        ("段階的リリース", "フェーズ分割でリスクを抑えて着実に。"),
    ]
    cw, gap, x0, y0, ch = 3.95, 0.32, 0.72, 1.95, 1.55
    for i, (title, body) in enumerate(cards):
        r, c = divmod(i, 3)
        x = x0 + c * (cw + gap)
        y = y0 + r * (ch + 0.28)
        rect(s, x, y, cw, ch, fill=WHITE, line=HAIR, radius=0.12, shadow=True)
        oval(s, x + 0.32, y + 0.32, 0.42, TINT)
        text(s, x + 0.32, y + 0.32, 0.42, 0.42, "✓", size=15, bold=True, color=BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + 0.92, y + 0.28, cw - 1.1, 0.4, title, size=14, bold=True, color=NAVY)
        text(s, x + 0.92, y + 0.72, cw - 1.15, 0.7, body, size=11, color=INK, line_spacing=1.25)
    cy = y0 + 2 * (ch + 0.28) + 0.02
    rect(s, x0, cy, (cw + gap) * 3 - gap, 0.9, fill=TINT, radius=0.12)
    oval(s, x0 + 0.32, cy + 0.28, 0.34, BLUE)
    text(s, x0 + 0.32, cy + 0.28, 0.34, 0.34, "★", size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x0 + 0.9, cy, (cw + gap) * 3 - gap - 1.1, 0.9,
         [("次のステップ：", {"size": 13, "bold": True, "color": NAVY}),
          ("まず要件定義フェーズからのスタートをおすすめします。仕様・スコープ・スケジュールを精緻化し、確実に前に進めます。",
           {"size": 12, "color": INK})],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
    footer(s)
    return s


def slide_closing(prs, meta):
    s = new(prs)
    bg(s, BLUE)
    text(s, SLIDE_W - 7.0, SLIDE_H - 2.6, 6.8, 2.0, "bravesoft", size=60, bold=True,
         color=RGBColor(0x3F, 0x86, 0xE6), align=PP_ALIGN.RIGHT)
    add_logo(s, SLIDE_W / 2 - 1.4, 2.15, 2.8, white=True)
    text(s, 1, 4.25, SLIDE_W - 2, 0.8, "ご清覧いただき、ありがとうございました。",
         size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sub = "{}　｜　{}".format(meta.get("client", "株式会社○○") + " 御中",
                             meta.get("project", "○○アプリ / システム開発") + " 御提案書")
    text(s, 1, 5.15, SLIDE_W - 2, 0.5, sub, size=12.5,
         color=RGBColor(0xD7, 0xE6, 0xFF), align=PP_ALIGN.CENTER)
    text(s, 1, 5.6, SLIDE_W - 2, 0.4,
         ("bravesoft株式会社　" + meta.get("date", "")).strip(),
         size=11, color=RGBColor(0xBF, 0xD6, 0xFF), align=PP_ALIGN.CENTER)
    return s


# ---------------------------------------------------------------- 組み立て
def generate(data, out_path):
    meta = data.get("meta", {})
    pats = compute(data)
    if not pats:
        raise ValueError("patterns（見積明細）が空です")
    primary = pats[0]

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide_cover(prs, meta)
    slide_points(prs, primary, meta)
    section_divider(prs, "01", "安心の根拠", "なぜ自信を持ってご提案できるか。")
    slide_reasons(prs)
    slide_donut(prs, primary)
    section_divider(prs, "02", "ご提案スコープ", "今回のお見積もりに含む機能の全体像。")
    slide_scope(prs, primary)
    section_divider(prs, "03", "お見積もり", "進め方と費用。金額はすべて税抜表記です。")
    slide_approach(prs, pats)
    slide_breakdown(prs, primary)
    slide_effort_bar(prs, primary)
    section_divider(prs, "04", "概算スケジュール", "フェーズ別の標準スケジュール。")
    slide_gantt(prs, primary)
    slide_summary(prs)
    slide_closing(prs, meta)

    prs.save(out_path)
    return out_path


def main():
    ap = argparse.ArgumentParser(description="bravesoft 提案書 PPTX 生成（Step 08）")
    ap.add_argument("--input", required=True, help="入力 JSON（Excel 生成と共通）")
    ap.add_argument("--out", required=True, help="出力 .pptx パス")
    args = ap.parse_args()
    with open(args.input, encoding="utf-8") as f:
        data = json.load(f)
    path = generate(data, args.out)
    print("OK: " + path)
    print("※ この提案書は提出用（社内原価・利益率は載せない＝Rule 10）。表紙の宛先・案件名は meta に従う。")


if __name__ == "__main__":
    main()
